"""대한조선 강의자료(docx/pptx/yml) → 텍스트(.md) 추출 → 지식 베이스에 저장.

RAG 인덱싱 전처리. docx/pptx는 임시 라이브러리로 텍스트만 뽑는다.
실행: uv run --with python-docx --with python-pptx python -m scripts.extract_lecture_docs
출력: data/knowledge/대한조선강의자료/<원본명>.md
"""
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent.parent
SRC = ROOT / "대한조선_강의자료"
OUT = ROOT / "data" / "knowledge" / "대한조선강의자료"


def from_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for tbl in doc.tables:  # 표 셀 텍스트도 포함
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def from_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"## 슬라이드 {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    n_ok = n_skip = 0
    for path in sorted(SRC.iterdir()):
        ext = path.suffix.lower()
        try:
            if ext == ".docx":
                text = from_docx(path)
            elif ext == ".pptx":
                text = from_pptx(path)
            elif ext in (".yml", ".yaml", ".txt", ".md"):
                text = path.read_text(encoding="utf-8", errors="ignore")
            else:
                continue
        except Exception as e:
            print(f"[ERR] {path.name}: {e}")
            continue
        if not text.strip():
            n_skip += 1
            print(f"[빈문서] {path.name}")
            continue
        dest = OUT / (path.stem + ".md")
        dest.write_text(f"# {path.name}\n\n{text}\n", encoding="utf-8")
        n_ok += 1
        print(f"[OK] {path.name} -> {len(text)}자")
    print(f"\n추출: {n_ok}개 저장, {n_skip}개 빈문서 → {OUT}")


if __name__ == "__main__":
    main()
